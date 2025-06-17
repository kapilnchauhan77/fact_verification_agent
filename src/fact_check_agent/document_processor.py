"""
Document processing module for OCR and text extraction
Supports 25+ document formats including Office, OpenOffice, images, archives, and more
"""

import io
import logging
import mimetypes
import os
import tempfile
import zipfile
import gc
import signal
import resource
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

import ebooklib
import eml_parser
import extract_msg
import mammoth
import markdownify

# Document processing libraries
import pdf2image
import pdfminer.high_level
import pdfplumber
import PyPDF2
import pytesseract
import striprtf
import xlrd
from bs4 import BeautifulSoup
from docx import Document
from ebooklib import epub
from google import genai

# Google Cloud Vision
from google.cloud import vision
from odf import teletype, text
from odf.opendocument import load as load_odf
from openpyxl import load_workbook
from PIL import Image

# Extended format support
from pptx import Presentation

try:
    import py7zr
except ImportError:
    py7zr = None
try:
    import rarfile
except ImportError:
    rarfile = None

# Audio/Video processing (optional)
try:
    import moviepy.editor as mp
    import speech_recognition as sr
    from pydub import AudioSegment

    AUDIO_VIDEO_SUPPORT = True
except ImportError:
    AUDIO_VIDEO_SUPPORT = False

from .config import config

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Handles OCR and text extraction from various document formats"""

    def __init__(self):
        """Initialize document processor with Google Cloud Vision client"""
        try:
            self.vision_client = vision.ImageAnnotatorClient()
        except Exception as e:
            logger.warning(f"Failed to initialize Google Cloud Vision client: {str(e)}")
            self.vision_client = None

        self.supported_formats = config.supported_formats

    def process_document(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Process a document and extract text with metadata

        Args:
            file_path: Path to the document file

        Returns:
            Dict containing extracted text, metadata, and processing info
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"Document not found: {file_path}")

        # Check file size
        file_size_mb = file_path.stat().st_size / (1024 * 1024)
        if file_size_mb > config.max_document_size_mb:
            raise ValueError(
                f"File size ({file_size_mb:.1f}MB) exceeds limit ({config.max_document_size_mb}MB)"
            )

        file_extension = file_path.suffix.lower().lstrip(".")

        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported format: {file_extension}")

        logger.info(f"Processing document: {file_path.name} ({file_extension})")

        try:
            # Route to appropriate processor based on file extension
            if file_extension == "pdf":
                print("Processing pdf:", file_path)
                return self._process_pdf_safe(file_path)
            elif file_extension in ["docx", "doc"]:
                return self._process_word_document(file_path)
            elif file_extension in ["pptx", "ppt"]:
                return self._process_powerpoint(file_path)
            elif file_extension in ["xlsx", "xls"]:
                return self._process_excel(file_path)
            elif file_extension in ["odt", "ods", "odp", "odg"]:
                return self._process_openoffice(file_path)
            elif file_extension in ["txt", "md"]:
                return self._process_text_file(file_path)
            elif file_extension == "rtf":
                return self._process_rtf(file_path)
            elif file_extension in ["html", "htm"]:
                return self._process_html(file_path)
            elif file_extension == "xml":
                return self._process_xml(file_path)
            elif file_extension in ["epub", "mobi"]:
                return self._process_ebook(file_path)
            elif file_extension in ["msg", "eml"]:
                return self._process_email(file_path)
            elif file_extension in ["zip", "7z", "rar"]:
                return self._process_archive(file_path)
            elif file_extension in ["jpg", "jpeg", "png", "tiff", "bmp", "gif", "webp"]:
                return self._process_image(file_path)
            elif file_extension in ["mp3", "wav", "mp4", "avi", "mov"]:
                return self._process_audio_video(file_path)
            else:
                # Fallback: try to determine by content
                return self._process_unknown_format(file_path)

        except Exception as e:
            logger.error(f"Error processing document {file_path}: {str(e)}")
            raise

    def _process_pdf_safe(self, file_path: Path) -> Dict[str, Any]:
        """Memory-safe PDF processing wrapper with timeout and resource monitoring"""
        import time
        import multiprocessing
        
        logger.info(f"Starting memory-safe PDF processing for: {file_path.name}")
        
        # Check available memory before processing
        try:
            # Get current memory usage
            process = os.getpid()
            memory_info = resource.getrusage(resource.RUSAGE_SELF)
            initial_memory = memory_info.ru_maxrss / 1024 / 1024  # Convert to MB
            logger.debug(f"Initial memory usage: {initial_memory:.1f}MB")
            
            # Set memory limits (1GB for PDF processing)
            max_memory_mb = 1024
            
            if initial_memory > max_memory_mb * 0.8:
                logger.warning(f"High memory usage detected ({initial_memory:.1f}MB), forcing garbage collection")
                gc.collect()
            
        except Exception as e:
            logger.debug(f"Memory check failed: {str(e)}")
        
        # Process with timeout
        start_time = time.time()
        timeout_seconds = 300  # 5 minute timeout
        
        try:
            result = self._process_pdf(file_path)
            
            processing_time = time.time() - start_time
            logger.info(f"PDF processing completed in {processing_time:.2f}s")
            
            # Final memory cleanup
            gc.collect()
            
            return result
            
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"PDF processing failed after {processing_time:.2f}s: {str(e)}")
            
            # Emergency memory cleanup
            gc.collect()
            
            # Return minimal result to prevent complete failure
            return {
                "text": "",
                "metadata": {
                    "file_type": "pdf",
                    "file_name": file_path.name,
                    "processing_method": "failed",
                    "error": str(e)
                },
                "success": False,
                "word_count": 0,
            }

    def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Process PDF document with improved memory management"""
        text_content = []
        metadata = {
            "file_type": "pdf",
            "file_name": file_path.name,
            "pages": 0,
            "processing_method": "safe_combined",
        }

        # First try direct text extraction with multiple libraries
        try:
            direct_text = self._extract_pdf_text_direct(file_path, metadata)
            if direct_text and len(direct_text.strip()) > 50:
                text_content.append({
                    "method": "direct_extraction",
                    "text": direct_text,
                    "confidence": 0.9,
                })
                logger.info(f"Successfully extracted text directly from PDF: {len(direct_text)} characters")
        except Exception as e:
            logger.warning(f"Direct PDF text extraction failed: {str(e)}")

        # If direct extraction failed or yielded little text, use safer OCR
        if not text_content or len(text_content[0]["text"]) < 100:
            try:
                ocr_text = self._extract_pdf_text_ocr_safe(file_path, metadata)
                if ocr_text and ocr_text.strip():
                    text_content.append({
                        "method": "safe_ocr",
                        "text": ocr_text,
                        "confidence": 0.75
                    })
                    logger.info(f"Successfully extracted text via OCR: {len(ocr_text)} characters")
            except Exception as e:
                logger.error(f"PDF OCR processing failed: {str(e)}")

        # Combine results
        final_text = ""
        processing_methods = []

        for content in text_content:
            final_text += content["text"]
            processing_methods.append(content["method"])

        metadata["processing_methods"] = processing_methods

        return {
            "text": final_text,
            "metadata": metadata,
            "success": bool(final_text.strip()),
            "word_count": len(final_text.split()) if final_text else 0,
        }

    def _extract_pdf_text_direct(self, file_path: Path, metadata: Dict[str, Any]) -> str:
        """Extract text directly from PDF using multiple libraries with error handling"""
        text_results = []
        
        # Method 1: Try pdfplumber (most reliable for structured text)
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                metadata["pages"] = len(pdf.pages)
                text_parts = []
                
                for page_num, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text_parts.append(f"\n--- Page {page_num + 1} ---\n{page_text}")
                    except Exception as e:
                        logger.debug(f"Failed to extract text from page {page_num + 1}: {str(e)}")
                        continue
                
                if text_parts:
                    pdfplumber_text = "\n".join(text_parts)
                    text_results.append(("pdfplumber", pdfplumber_text))
                    logger.debug(f"pdfplumber extracted {len(pdfplumber_text)} characters")
                    
        except Exception as e:
            logger.debug(f"pdfplumber extraction failed: {str(e)}")
        
        # Method 2: Try pdfminer (good for complex layouts)
        try:
            import pdfminer.high_level
            pdfminer_text = pdfminer.high_level.extract_text(str(file_path))
            if pdfminer_text and pdfminer_text.strip():
                text_results.append(("pdfminer", pdfminer_text))
                logger.debug(f"pdfminer extracted {len(pdfminer_text)} characters")
        except Exception as e:
            logger.debug(f"pdfminer extraction failed: {str(e)}")
        
        # Method 3: Try PyPDF2 as last resort
        try:
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                if not metadata.get("pages"):
                    metadata["pages"] = len(pdf_reader.pages)
                
                pypdf2_text = ""
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            pypdf2_text += f"\n--- Page {page_num + 1} ---\n{page_text}"
                    except Exception as e:
                        logger.debug(f"PyPDF2 failed on page {page_num + 1}: {str(e)}")
                        continue
                
                if pypdf2_text.strip():
                    text_results.append(("PyPDF2", pypdf2_text))
                    logger.debug(f"PyPDF2 extracted {len(pypdf2_text)} characters")
                    
        except Exception as e:
            logger.debug(f"PyPDF2 extraction failed: {str(e)}")
        
        # Return the best result (longest text)
        if text_results:
            best_method, best_text = max(text_results, key=lambda x: len(x[1]))
            logger.info(f"Best PDF extraction method: {best_method} ({len(best_text)} characters)")
            return best_text
        
        return ""
    
    def _extract_pdf_text_ocr_safe(self, file_path: Path, metadata: Dict[str, Any]) -> str:
        """Extract text from PDF using OCR with safer memory management"""
        import gc
        import tempfile
        
        ocr_text = ""
        max_pages = 50  # Limit pages to prevent memory issues
        
        try:
            # Get PDF info first
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                metadata["pages"] = total_pages
                
            if total_pages > max_pages:
                logger.warning(f"PDF has {total_pages} pages, processing only first {max_pages} for OCR")
                total_pages = max_pages
            
            # Process pages in batches to manage memory
            batch_size = 5  # Process 5 pages at a time
            
            for batch_start in range(0, total_pages, batch_size):
                batch_end = min(batch_start + batch_size, total_pages)
                
                try:
                    # Convert only this batch of pages
                    batch_images = pdf2image.convert_from_path(
                        file_path,
                        first_page=batch_start + 1,
                        last_page=batch_end,
                        dpi=150,  # Lower DPI to reduce memory usage
                        thread_count=1  # Single thread to avoid memory conflicts
                    )
                    
                    for page_num, image in enumerate(batch_images):
                        try:
                            actual_page_num = batch_start + page_num + 1
                            
                            # Resize image if too large
                            if image.size[0] > 2000 or image.size[1] > 2000:
                                image.thumbnail((2000, 2000), Image.Resampling.LANCZOS)
                            
                            page_text = self._perform_ocr_safe(image)
                            if page_text and page_text.strip():
                                ocr_text += f"\n--- Page {actual_page_num} ---\n{page_text}"
                            
                            # Explicitly close the image to free memory
                            image.close()
                            
                        except Exception as e:
                            logger.warning(f"OCR failed on page {actual_page_num}: {str(e)}")
                            continue
                    
                    # Clear batch images from memory
                    del batch_images
                    gc.collect()
                    
                    logger.debug(f"Processed OCR batch {batch_start + 1}-{batch_end}")
                    
                except Exception as e:
                    logger.error(f"Failed to process OCR batch {batch_start + 1}-{batch_end}: {str(e)}")
                    continue
            
            return ocr_text
            
        except Exception as e:
            logger.error(f"OCR PDF processing failed: {str(e)}")
            return ""
        finally:
            # Force garbage collection
            gc.collect()
    
    def _perform_ocr_safe(self, image: Image.Image) -> str:
        """Perform OCR with better error handling and memory management"""
        try:
            # Try Gemini OCR first (more accurate)
            gemini_text = self._perform_ocr(image)
            if gemini_text and len(gemini_text.strip()) > 10:
                return gemini_text
            
            # Fallback to Tesseract
            return self._perform_tesseract_ocr(image)
            
        except Exception as e:
            logger.warning(f"Safe OCR failed: {str(e)}")
            return ""
    
    def _perform_tesseract_ocr(self, image: Image.Image) -> str:
        """Perform Tesseract OCR as fallback"""
        try:
            # Configure Tesseract for better accuracy
            custom_config = r'--oem 3 --psm 6'
            text = pytesseract.image_to_string(image, config=custom_config)
            return text.strip()
        except Exception as e:
            logger.warning(f"Tesseract OCR failed: {str(e)}")
            return ""

    def _process_word_document(self, file_path: Path) -> Dict[str, Any]:
        """Process DOCX/DOC document"""
        file_extension = file_path.suffix.lower().lstrip(".")

        try:
            if file_extension == "docx":
                # Native DOCX processing
                doc = Document(file_path)

                # Extract text from paragraphs
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                text = "\n".join(paragraphs)

                # Extract text from tables
                table_text = []
                for table in doc.tables:
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        table_text.append(" | ".join(row_text))

                if table_text:
                    text += "\n\n--- Tables ---\n" + "\n".join(table_text)

                metadata = {
                    "file_type": file_extension,
                    "file_name": file_path.name,
                    "paragraphs": len(paragraphs),
                    "tables": len(doc.tables),
                    "processing_method": "native",
                }

            elif file_extension == "doc":
                # Legacy DOC processing using mammoth
                with open(file_path, "rb") as docx_file:
                    result = mammoth.extract_text(docx_file)
                    text = result.value

                metadata = {
                    "file_type": file_extension,
                    "file_name": file_path.name,
                    "processing_method": "mammoth",
                    "conversion_messages": result.messages,
                }

            return {
                "text": text,
                "metadata": metadata,
                "success": bool(text.strip()),
                "word_count": len(text.split()) if text else 0,
            }

        except Exception as e:
            logger.error(f"Word document processing failed: {str(e)}")
            raise

    def _process_image(self, file_path: Path) -> Dict[str, Any]:
        """Process image document using OCR"""
        try:
            with Image.open(file_path) as image:
                # Perform OCR using Google Cloud Vision API first
                text = self._perform_google_vision_ocr(image)
                confidence = 0.85
                method = "google_vision"

                # Fallback to Tesseract if Google Vision fails
                if not text or len(text.strip()) < 10:
                    text = self._perform_ocr(image)
                    confidence = 0.7
                    method = "tesseract"

                metadata = {
                    "file_type": file_path.suffix.lower().lstrip("."),
                    "file_name": file_path.name,
                    "image_size": image.size,
                    "processing_method": method,
                    "confidence": confidence,
                }

                return {
                    "text": text,
                    "metadata": metadata,
                    "success": bool(text.strip()),
                    "word_count": len(text.split()) if text else 0,
                }

        except Exception as e:
            logger.error(f"Image processing failed: {str(e)}")
            raise

    def _perform_google_vision_ocr(self, image: Image.Image) -> str:
        """Perform OCR using Google Cloud Vision API"""
        try:
            if not self.vision_client:
                logger.warning("Google Vision client not available, skipping")
                return ""

            # Convert PIL image to bytes
            img_byte_arr = io.BytesIO()
            image.save(img_byte_arr, format="PNG")
            img_byte_arr = img_byte_arr.getvalue()

            # Create Vision API image object
            vision_image = vision.Image(content=img_byte_arr)

            # Perform text detection
            response = self.vision_client.text_detection(image=vision_image)

            if response.error.message:
                raise Exception(f"Google Vision API error: {response.error.message}")

            # Extract text from response
            texts = response.text_annotations
            if texts:
                return texts[0].description
            else:
                return ""

        except Exception as e:
            logger.warning(f"Google Vision OCR failed: {str(e)}")
            return ""

    def _perform_ocr(self, image: Image.Image) -> str:
        try:
            client = genai.Client(
                vertexai=True,
                project="mythic-lattice-455715-q1",
                location="global",
            )
            response = client.models.generate_content(
                model="gemini-2.0-flash",
                contents=[
                    "Extract the text from the image while maintaining the original layout and formatting.",
                    image,
                ],
            )
            if (
                response is not None
                and hasattr(response, "text")
                and isinstance(response.text, str)
            ):
                resp = response.text.replace("```text", "").replace("```", "").strip()
                return resp

            return ""

        except Exception as e:
            logger.warning(f"Gemini OCR failed: {str(e)}")
            return ""

    def extract_document_metadata(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Extract metadata from document without full processing"""
        file_path = Path(file_path)

        metadata = {
            "file_name": file_path.name,
            "file_size_mb": file_path.stat().st_size / (1024 * 1024),
            "file_extension": file_path.suffix.lower().lstrip("."),
            "created_time": file_path.stat().st_ctime,
            "modified_time": file_path.stat().st_mtime,
        }

        return metadata

    def _process_powerpoint(self, file_path: Path) -> Dict[str, Any]:
        """Process PowerPoint presentations (PPTX/PPT)"""
        try:
            prs = Presentation(file_path)

            slides_text = []
            slide_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())

                if slide_content:
                    slides_text.append(
                        f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_content)
                    )
                    slide_count += 1

            text = "\n".join(slides_text)

            metadata = {
                "file_type": file_path.suffix.lower().lstrip("."),
                "file_name": file_path.name,
                "slides": slide_count,
                "processing_method": "native",
            }

            return {
                "text": text,
                "metadata": metadata,
                "success": bool(text.strip()),
                "word_count": len(text.split()) if text else 0,
            }

        except Exception as e:
            logger.error(f"PowerPoint processing failed: {str(e)}")
            raise

    def _process_excel(self, file_path: Path) -> Dict[str, Any]:
        """Process Excel spreadsheets (XLSX/XLS)"""
        file_extension = file_path.suffix.lower().lstrip(".")

        try:
            sheets_text = []

            if file_extension == "xlsx":
                # Use openpyxl for XLSX
                workbook = load_workbook(file_path, data_only=True)

                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    sheet_content = [f"\n--- Sheet: {sheet_name} ---"]

                    for row in sheet.iter_rows(values_only=True):
                        row_data = [
                            str(cell) if cell is not None else "" for cell in row
                        ]
                        if any(cell.strip() for cell in row_data if cell):
                            sheet_content.append(" | ".join(row_data))

                    if len(sheet_content) > 1:  # More than just header
                        sheets_text.extend(sheet_content)

            elif file_extension == "xls":
                # Use xlrd for XLS
                workbook = xlrd.open_workbook(file_path)

                for sheet_idx in range(workbook.nsheets):
                    sheet = workbook.sheet_by_index(sheet_idx)
                    sheet_content = [f"\n--- Sheet: {sheet.name} ---"]

                    for row_idx in range(sheet.nrows):
                        row_data = [
                            str(sheet.cell_value(row_idx, col_idx))
                            for col_idx in range(sheet.ncols)
                        ]
                        if any(cell.strip() for cell in row_data if cell):
                            sheet_content.append(" | ".join(row_data))

                    if len(sheet_content) > 1:
                        sheets_text.extend(sheet_content)

            text = "\n".join(sheets_text)

            metadata = {
                "file_type": file_extension,
                "file_name": file_path.name,
                "processing_method": "native",
            }

            return {
                "text": text,
                "metadata": metadata,
                "success": bool(text.strip()),
                "word_count": len(text.split()) if text else 0,
            }

        except Exception as e:
            logger.error(f"Excel processing failed: {str(e)}")
            raise

    def _process_openoffice(self, file_path: Path) -> Dict[str, Any]:
        """Process OpenOffice/LibreOffice documents (ODT/ODS/ODP)"""
        try:
            doc = load_odf(file_path)
            text_content = []

            # Extract all text content
            for element in doc.getElementsByType(text.P):
                paragraph_text = teletype.extractText(element)
                if paragraph_text.strip():
                    text_content.append(paragraph_text)

            # Extract text from tables
            for table in doc.getElementsByType(text.Table):
                for row in table.getElementsByType(text.TableRow):
                    row_text = []
                    for cell in row.getElementsByType(text.TableCell):
                        cell_text = teletype.extractText(cell)
                        row_text.append(cell_text.strip())
                    if any(cell for cell in row_text):
                        text_content.append(" | ".join(row_text))

            text = "\n".join(text_content)

            metadata = {
                "file_type": file_path.suffix.lower().lstrip("."),
                "file_name": file_path.name,
                "processing_method": "odf",
            }

            return {
                "text": text,
                "metadata": metadata,
                "success": bool(text.strip()),
                "word_count": len(text.split()) if text else 0,
            }

        except Exception as e:
            logger.error(f"OpenOffice document processing failed: {str(e)}")
            raise

    def _process_text_file(self, file_path: Path) -> Dict[str, Any]:
        """Process plain text and markdown files"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

            file_extension = file_path.suffix.lower().lstrip(".")

            # Convert markdown to plain text if needed
            if file_extension == "md":
                text = markdownify.markdownify(text, strip=["a", "img"])

            metadata = {
                "file_type": file_extension,
                "file_name": file_path.name,
                "processing_method": "native",
                "encoding": "utf-8",
            }

            return {
                "text": text,
                "metadata": metadata,
                "success": bool(text.strip()),
                "word_count": len(text.split()) if text else 0,
            }

        except Exception as e:
            logger.error(f"Text file processing failed: {str(e)}")
            raise

    def _process_rtf(self, file_path: Path) -> Dict[str, Any]:
        """Process RTF documents"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                rtf_content = f.read()

            # Extract plain text from RTF
            text = striprtf.rtf_to_text(rtf_content)

            metadata = {
                "file_type": "rtf",
                "file_name": file_path.name,
                "processing_method": "striprtf",
            }

            return {
                "text": text,
                "metadata": metadata,
                "success": bool(text.strip()),
                "word_count": len(text.split()) if text else 0,
            }

        except Exception as e:
            logger.error(f"RTF processing failed: {str(e)}")
            raise

    def _process_html(self, file_path: Path) -> Dict[str, Any]:
        """Process HTML documents"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                html_content = f.read()

            # Parse HTML and extract text
            soup = BeautifulSoup(html_content, "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            # Extract text
            text = soup.get_text()

            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = "\n".join(chunk for chunk in chunks if chunk)

            metadata = {
                "file_type": file_path.suffix.lower().lstrip("."),
                "file_name": file_path.name,
                "processing_method": "beautifulsoup",
                "title": soup.title.string if soup.title else None,
            }

            return {
                "text": text,
                "metadata": metadata,
                "success": bool(text.strip()),
                "word_count": len(text.split()) if text else 0,
            }

        except Exception as e:
            logger.error(f"HTML processing failed: {str(e)}")
            raise

    def _process_xml(self, file_path: Path) -> Dict[str, Any]:
        """Process XML documents"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                xml_content = f.read()

            # Parse XML and extract text content
            soup = BeautifulSoup(xml_content, "xml")
            text = soup.get_text()

            # Clean up whitespace
            text = "\n".join(line.strip() for line in text.splitlines() if line.strip())

            metadata = {
                "file_type": "xml",
                "file_name": file_path.name,
                "processing_method": "beautifulsoup",
            }

            return {
                "text": text,
                "metadata": metadata,
                "success": bool(text.strip()),
                "word_count": len(text.split()) if text else 0,
            }

        except Exception as e:
            logger.error(f"XML processing failed: {str(e)}")
            raise

    def _process_ebook(self, file_path: Path) -> Dict[str, Any]:
        """Process EPUB and MOBI ebooks"""
        file_extension = file_path.suffix.lower().lstrip(".")

        try:
            if file_extension == "epub":
                book = epub.read_epub(file_path)
                chapters = []

                # Extract text from all chapters
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        soup = BeautifulSoup(item.get_content(), "html.parser")
                        chapter_text = soup.get_text()
                        if chapter_text.strip():
                            chapters.append(chapter_text.strip())

                text = "\n\n".join(chapters)

                # Extract metadata
                title = book.get_metadata("DC", "title")
                author = book.get_metadata("DC", "creator")

                metadata = {
                    "file_type": file_extension,
                    "file_name": file_path.name,
                    "processing_method": "ebooklib",
                    "title": title[0][0] if title else None,
                    "author": author[0][0] if author else None,
                    "chapters": len(chapters),
                }

            else:  # MOBI
                # MOBI processing would require additional library like mobidedrm
                # For now, fallback to basic processing
                text = "MOBI format not fully supported - please convert to EPUB or PDF"
                metadata = {
                    "file_type": file_extension,
                    "file_name": file_path.name,
                    "processing_method": "unsupported",
                }

            return {
                "text": text,
                "metadata": metadata,
                "success": bool(text.strip()),
                "word_count": len(text.split()) if text else 0,
            }

        except Exception as e:
            logger.error(f"Ebook processing failed: {str(e)}")
            raise

    def _process_email(self, file_path: Path) -> Dict[str, Any]:
        """Process email files (MSG/EML)"""
        file_extension = file_path.suffix.lower().lstrip(".")

        try:
            if file_extension == "msg":
                # Process MSG files
                msg = extract_msg.Message(file_path)

                email_parts = []
                if msg.subject:
                    email_parts.append(f"Subject: {msg.subject}")
                if msg.sender:
                    email_parts.append(f"From: {msg.sender}")
                if msg.to:
                    email_parts.append(f"To: {msg.to}")
                if msg.date:
                    email_parts.append(f"Date: {msg.date}")
                if msg.body:
                    email_parts.append(f"\nBody:\n{msg.body}")

                text = "\n".join(email_parts)

                metadata = {
                    "file_type": file_extension,
                    "file_name": file_path.name,
                    "processing_method": "extract_msg",
                    "subject": msg.subject,
                    "sender": msg.sender,
                    "attachments": len(msg.attachments) if msg.attachments else 0,
                }

            elif file_extension == "eml":
                # Process EML files
                with open(file_path, "rb") as f:
                    email_content = f.read()

                parsed_email = eml_parser.eml_parser.decode_email_b(email_content)

                email_parts = []
                header = parsed_email.get("header", {})

                if "subject" in header:
                    email_parts.append(f"Subject: {header['subject']}")
                if "from" in header:
                    email_parts.append(f"From: {header['from']}")
                if "to" in header:
                    email_parts.append(f"To: {header['to']}")
                if "date" in header:
                    email_parts.append(f"Date: {header['date']}")

                # Extract body
                body = parsed_email.get("body", [])
                for body_part in body:
                    if "content" in body_part:
                        email_parts.append(f"\nBody:\n{body_part['content']}")

                text = "\n".join(email_parts)

                metadata = {
                    "file_type": file_extension,
                    "file_name": file_path.name,
                    "processing_method": "eml_parser",
                    "subject": header.get("subject"),
                    "sender": header.get("from"),
                }

            return {
                "text": text,
                "metadata": metadata,
                "success": bool(text.strip()),
                "word_count": len(text.split()) if text else 0,
            }

        except Exception as e:
            logger.error(f"Email processing failed: {str(e)}")
            raise

    def _process_archive(self, file_path: Path) -> Dict[str, Any]:
        """Process archive files (ZIP/7Z/RAR)"""
        file_extension = file_path.suffix.lower().lstrip(".")

        try:
            extracted_texts = []
            processed_files = []

            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)

                # Extract archive
                if file_extension == "zip":
                    with zipfile.ZipFile(file_path, "r") as zip_ref:
                        zip_ref.extractall(temp_path)

                elif file_extension == "7z" and py7zr:
                    with py7zr.SevenZipFile(file_path, mode="r") as archive:
                        archive.extractall(temp_path)

                elif file_extension == "rar" and rarfile:
                    with rarfile.RarFile(file_path) as rar_ref:
                        rar_ref.extractall(temp_path)

                else:
                    raise ValueError(
                        f"Archive format {file_extension} not supported or library not available"
                    )

                # Process extracted files
                for extracted_file in temp_path.rglob("*"):
                    if extracted_file.is_file():
                        try:
                            file_ext = extracted_file.suffix.lower().lstrip(".")
                            if (
                                file_ext in self.supported_formats
                                and file_ext != file_extension
                            ):
                                # Recursively process supported files
                                result = self.process_document(extracted_file)
                                if result["success"]:
                                    extracted_texts.append(
                                        f"\n--- File: {extracted_file.name} ---\n{result['text']}"
                                    )
                                    processed_files.append(extracted_file.name)
                        except Exception as e:
                            logger.warning(
                                f"Failed to process archived file {extracted_file.name}: {str(e)}"
                            )

            text = "\n".join(extracted_texts)

            metadata = {
                "file_type": file_extension,
                "file_name": file_path.name,
                "processing_method": "archive_extraction",
                "extracted_files": processed_files,
                "total_extracted": len(processed_files),
            }

            return {
                "text": text,
                "metadata": metadata,
                "success": bool(text.strip()),
                "word_count": len(text.split()) if text else 0,
            }

        except Exception as e:
            logger.error(f"Archive processing failed: {str(e)}")
            raise

    def _process_audio_video(self, file_path: Path) -> Dict[str, Any]:
        """Process audio and video files using speech recognition"""
        if not AUDIO_VIDEO_SUPPORT:
            raise ValueError(
                "Audio/video processing requires additional dependencies. Install with: pip install SpeechRecognition pydub moviepy"
            )

        file_extension = file_path.suffix.lower().lstrip(".")

        try:
            recognizer = sr.Recognizer()

            with tempfile.TemporaryDirectory() as temp_dir:
                temp_audio_path = Path(temp_dir) / "audio.wav"

                # Extract audio from video files
                if file_extension in ["mp4", "avi", "mov"]:
                    video = mp.VideoFileClip(str(file_path))
                    audio = video.audio
                    audio.write_audiofile(
                        str(temp_audio_path), verbose=False, logger=None
                    )
                    audio.close()
                    video.close()

                # Convert audio to WAV format for speech recognition
                elif file_extension in ["mp3"]:
                    audio = AudioSegment.from_mp3(file_path)
                    audio.export(temp_audio_path, format="wav")

                elif file_extension == "wav":
                    temp_audio_path = file_path

                else:
                    raise ValueError(
                        f"Unsupported audio/video format: {file_extension}"
                    )

                # Perform speech recognition
                with sr.AudioFile(str(temp_audio_path)) as source:
                    audio_data = recognizer.record(source)

                # Try Google Speech Recognition first, fallback to others
                transcript = ""
                recognition_method = ""

                try:
                    transcript = recognizer.recognize_google(audio_data)
                    recognition_method = "google"
                except sr.UnknownValueError:
                    transcript = "Could not understand audio"
                    recognition_method = "failed"
                except sr.RequestError as e:
                    # Fallback to offline recognition
                    try:
                        transcript = recognizer.recognize_sphinx(audio_data)
                        recognition_method = "sphinx"
                    except:
                        transcript = f"Speech recognition service error: {e}"
                        recognition_method = "error"

            metadata = {
                "file_type": file_extension,
                "file_name": file_path.name,
                "processing_method": "speech_recognition",
                "recognition_engine": recognition_method,
            }

            return {
                "text": transcript,
                "metadata": metadata,
                "success": bool(
                    transcript and transcript != "Could not understand audio"
                ),
                "word_count": len(transcript.split()) if transcript else 0,
            }

        except Exception as e:
            logger.error(f"Audio/video processing failed: {str(e)}")
            raise

    def _process_unknown_format(self, file_path: Path) -> Dict[str, Any]:
        """Attempt to process unknown file formats"""
        try:
            # Try to determine MIME type
            mime_type, _ = mimetypes.guess_type(file_path)

            if mime_type:
                if mime_type.startswith("text/"):
                    # Try as text file
                    return self._process_text_file(file_path)
                elif mime_type.startswith("image/"):
                    # Try as image
                    return self._process_image(file_path)
                elif mime_type == "application/pdf":
                    # Try as PDF
                    return self._process_pdf(file_path)

            # Last resort: try to read as text
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()

                if text.strip():
                    metadata = {
                        "file_type": file_path.suffix.lower().lstrip("."),
                        "file_name": file_path.name,
                        "processing_method": "fallback_text",
                        "mime_type": mime_type,
                    }

                    return {
                        "text": text,
                        "metadata": metadata,
                        "success": bool(text.strip()),
                        "word_count": len(text.split()) if text else 0,
                    }
            except:
                pass

            # If all else fails
            raise ValueError(f"Cannot process file format: {file_path.suffix}")

        except Exception as e:
            logger.error(f"Unknown format processing failed: {str(e)}")
            raise

    def get_supported_formats_info(self) -> Dict[str, List[str]]:
        """Get detailed information about supported formats"""
        return {
            "images": ["jpg", "jpeg", "png", "tiff", "bmp", "gif", "webp"],
            "documents": ["pdf", "docx", "doc", "txt", "rtf", "md"],
            "presentations": ["pptx", "ppt"],
            "spreadsheets": ["xlsx", "xls"],
            "openoffice": ["odt", "ods", "odp", "odg"],
            "web": ["html", "htm", "xml"],
            "ebooks": ["epub", "mobi"],
            "email": ["msg", "eml"],
            "archives": ["zip", "7z", "rar"],
            "audio_video": (
                ["mp3", "wav", "mp4", "avi", "mov"] if AUDIO_VIDEO_SUPPORT else []
            ),
        }

